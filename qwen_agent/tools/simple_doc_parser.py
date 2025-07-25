# Copyright 2023 The Qwen team, Alibaba Group. All rights reserved.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#    http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import json
import os
import re
import time
from collections import Counter
from typing import Dict, List, Optional, Union

from qwen_agent.log import logger
from qwen_agent.settings import DEFAULT_WORKSPACE
from qwen_agent.tools.base import BaseTool, register_tool
from qwen_agent.tools.storage import KeyNotExistsError, Storage
from qwen_agent.utils.str_processing import rm_cid, rm_continuous_placeholders, rm_hexadecimal
from qwen_agent.utils.tokenization_qwen import count_tokens
from qwen_agent.utils.utils import (get_file_type, hash_sha256, is_http_url, read_text_from_file,
                                    sanitize_chrome_file_path, save_url_to_local_work_dir)


def clean_paragraph(text):
    text = rm_cid(text)
    text = rm_hexadecimal(text)
    text = rm_continuous_placeholders(text)
    return text


class DocParserError(Exception):

    def __init__(self,
                 exception: Optional[Exception] = None,
                 code: Optional[str] = None,
                 message: Optional[str] = None,
                 extra: Optional[dict] = None):
        if exception is not None:
            super().__init__(exception)
        else:
            super().__init__(f'\nError code: {code}. Error message: {message}')
        self.exception = exception
        self.code = code
        self.message = message
        self.extra = extra


PARAGRAPH_SPLIT_SYMBOL = '\n'


def parse_word(docx_path: str, extract_image: bool = False):
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    from docx import Document
    doc = Document(docx_path)

    content = []
    for para in doc.paragraphs:
        content.append({'text': para.text})
    for table in doc.tables:
        tbl = []
        for row in table.rows:
            tbl.append('|' + '|'.join([cell.text for cell in row.cells]) + '|')
        tbl = '\n'.join(tbl)
        content.append({'table': tbl})

    # Due to the pages in Word are not fixed, the entire document is returned as one page
    return [{'page_num': 1, 'content': content}]


def parse_ppt(path: str, extract_image: bool = False):
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    try:
        ppt = Presentation(path)
    except PackageNotFoundError as ex:
        logger.warning(ex)
        return []
    doc = []
    for slide_number, slide in enumerate(ppt.slides):
        page = {'page_num': slide_number + 1, 'content': []}

        for shape in slide.shapes:
            if not shape.has_text_frame and not shape.has_table:
                pass

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = ''.join(run.text for run in paragraph.runs)
                    paragraph_text = clean_paragraph(paragraph_text)
                    if paragraph_text.strip():
                        page['content'].append({'text': paragraph_text})

            if shape.has_table:
                tbl = []
                for row_number, row in enumerate(shape.table.rows):
                    tbl.append('|' + '|'.join([cell.text for cell in row.cells]) + '|')
                tbl = '\n'.join(tbl)
                page['content'].append({'table': tbl})
        doc.append(page)
    return doc


def parse_txt(path: str):
    text = read_text_from_file(path)
    paras = text.split(PARAGRAPH_SPLIT_SYMBOL)
    content = []
    for p in paras:
        content.append({'text': p})

    # Due to the pages in txt are not fixed, the entire document is returned as one page
    return [{'page_num': 1, 'content': content}]


def df_to_md(df) -> str:

    def replace_long_dashes(text):
        if text.replace('-', '').replace(':', '').strip():
            return text
        pattern = r'-{6,}'
        replaced_text = re.sub(pattern, '-----', text)
        return replaced_text

    from tabulate import tabulate
    df = df.dropna(how='all')
    df = df.dropna(axis=1, how='all')
    df = df.fillna('')
    md_table = tabulate(df, headers='keys', tablefmt='pipe', showindex=False)

    md_table = '\n'.join([
        '|'.join(replace_long_dashes(' ' + cell.strip() + ' ' if cell else '')
                 for cell in row.split('|'))
        for row in md_table.split('\n')
    ])
    return md_table


def parse_excel(file_path: str, extract_image: bool = False) -> List[dict]:
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    import pandas as pd

    excel_file = pd.ExcelFile(file_path)
    md_tables = []
    for sheet_name in excel_file.sheet_names:
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        md_table = df_to_md(df)
        md_tables.append(f'### Sheet: {sheet_name}\n{md_table}')

    return [{'page_num': i + 1, 'content': [{'table': md_tables[i]}]} for i in range(len(md_tables))]


def parse_csv(file_path: str, extract_image: bool = False) -> List[dict]:
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    import pandas as pd
    md_tables = []
    try:
        df = pd.read_csv(file_path, encoding_errors='replace', on_bad_lines='skip')
    except Exception as ex:
        # Directly converted from Excel
        logger.warning(ex)
        return parse_excel(file_path, extract_image)
    md_table = df_to_md(df)
    md_tables.append(md_table)  # There is only one table available

    return [{'page_num': i + 1, 'content': [{'table': md_tables[i]}]} for i in range(len(md_tables))]


def parse_tsv(file_path: str, extract_image: bool = False) -> List[dict]:
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    import pandas as pd
    md_tables = []
    try:
        df = pd.read_csv(file_path, sep='\t', encoding_errors='replace', on_bad_lines='skip')
    except Exception as ex:
        # Directly converted from Excel
        logger.warning(ex)
        return parse_excel(file_path, extract_image)
    md_table = df_to_md(df)
    md_tables.append(md_table)  # There is only one table available

    return [{'page_num': i + 1, 'content': [{'table': md_tables[i]}]} for i in range(len(md_tables))]


def parse_html_bs(path: str, extract_image: bool = False):
    if extract_image:
        raise ValueError('Currently, extracting images is not supported!')

    def pre_process_html(s):
        # replace multiple newlines
        s = re.sub('\n+', '\n', s)
        # replace special string
        s = s.replace("Add to Qwen's Reading List", '')
        return s

    try:
        from bs4 import BeautifulSoup
    except Exception:
        raise ValueError('Please install bs4 by `pip install beautifulsoup4`')
    bs_kwargs = {'features': 'lxml'}
    with open(path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, **bs_kwargs)

    text = soup.get_text()

    if soup.title:
        title = str(soup.title.string)
    else:
        title = ''

    text = pre_process_html(text)
    paras = text.split(PARAGRAPH_SPLIT_SYMBOL)
    content = []
    for p in paras:
        p = clean_paragraph(p)
        if p.strip():
            content.append({'text': p})

    # The entire document is returned as one page
    return [{'page_num': 1, 'content': content, 'title': title}]


def parse_pdf(pdf_path: str, extract_image: bool = False) -> List[dict]:
    # Todo: header and footer
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTImage, LTRect, LTTextContainer

    doc = []
    import pdfplumber
    pdf = pdfplumber.open(pdf_path)
    for i, page_layout in enumerate(extract_pages(pdf_path)):
        page = {'page_num': page_layout.pageid, 'content': []}

        elements = []
        for element in page_layout:
            elements.append(element)

        # Init params for table
        table_num = 0
        tables = []

        for element in elements:
            if isinstance(element, LTRect):
                if not tables:
                    tables = extract_tables(pdf, i)
                if table_num < len(tables):
                    table_string = table_converter(tables[table_num])
                    table_num += 1
                    if table_string:
                        page['content'].append({'table': table_string, 'obj': element})
            elif isinstance(element, LTTextContainer):
                # Delete line breaks in the same paragraph
                text = element.get_text()
                # Todo: Further analysis using font
                font = get_font(element)
                if text.strip():
                    new_content_item = {'text': text, 'obj': element}
                    if font:
                        new_content_item['font-size'] = round(font[1])
                        # new_content_item['font-name'] = font[0]
                    page['content'].append(new_content_item)
            elif extract_image and isinstance(element, LTImage):
                # Todo: ocr
                raise ValueError('Currently, extracting images is not supported!')
            else:
                pass

        # merge elements
        page['content'] = postprocess_page_content(page['content'])
        doc.append(page)

    return doc


def postprocess_page_content(page_content: list) -> list:
    # rm repetitive identification for table and text
    # Some documents may repeatedly recognize LTRect and LTTextContainer
    table_obj = [p['obj'] for p in page_content if 'table' in p]
    tmp = []
    for p in page_content:
        repetitive = False
        if 'text' in p:
            for t in table_obj:
                if t.bbox[0] <= p['obj'].bbox[0] and p['obj'].bbox[1] <= t.bbox[1] and t.bbox[2] <= p['obj'].bbox[
                        2] and p['obj'].bbox[3] <= t.bbox[3]:
                    repetitive = True
                    break

        if not repetitive:
            tmp.append(p)
    page_content = tmp

    # merge paragraphs that have been separated by mistake
    new_page_content = []
    for p in page_content:
        if new_page_content and 'text' in new_page_content[-1] and 'text' in p and abs(
                p.get('font-size', 12) -
                new_page_content[-1].get('font-size', 12)) < 2 and p['obj'].height < p.get('font-size', 12) + 1:
            # Merge those lines belonging to a paragraph
            _p = p['text']
            new_page_content[-1]['text'] += f' {_p}'
            # new_page_content[-1]['font-name'] = p.get('font-name', '')
            new_page_content[-1]['font-size'] = p.get('font-size', 12)
        else:
            p.pop('obj')
            new_page_content.append(p)
    for i in range(len(new_page_content)):
        if 'text' in new_page_content[i]:
            new_page_content[i]['text'] = clean_paragraph(new_page_content[i]['text'])
    return new_page_content


def get_font(element):
    from pdfminer.layout import LTChar, LTTextContainer

    fonts_list = []
    for text_line in element:
        if isinstance(text_line, LTTextContainer):
            for character in text_line:
                if isinstance(character, LTChar):
                    fonts_list.append((character.fontname, character.size))

    fonts_list = list(set(fonts_list))
    if fonts_list:
        counter = Counter(fonts_list)
        most_common_fonts = counter.most_common(1)[0][0]
        return most_common_fonts
    else:
        return []


def extract_tables(pdf, page_num):
    table_page = pdf.pages[page_num]
    tables = table_page.extract_tables()
    return tables


def table_converter(table):
    table_string = ''
    for row_num in range(len(table)):
        row = table[row_num]
        cleaned_row = [
            item.replace('\n', ' ') if item is not None and '\n' in item else 'None' if item is None else item
            for item in row
        ]
        table_string += ('|' + '|'.join(cleaned_row) + '|' + '\n')
    table_string = table_string[:-1]
    return table_string


PARSER_SUPPORTED_FILE_TYPES = ['pdf', 'docx', 'pptx', 'txt', 'html', 'csv', 'tsv', 'xlsx', 'xls']


def get_plain_doc(doc: list):
    paras = []
    for page in doc:
        for para in page['content']:
            for k, v in para.items():
                if k in ['text', 'table', 'image']:
                    paras.append(v)
    return PARAGRAPH_SPLIT_SYMBOL.join(paras)


@register_tool('simple_doc_parser')
class SimpleDocParser(BaseTool):
    description = f"提取出一个文档的内容，支持类型包括：{' / '.join(PARSER_SUPPORTED_FILE_TYPES)}"
    parameters = {
        'type': 'object',
        'properties': {
            'url': {
                'description': '待提取的文件的路径，可以是一个本地路径或可下载的http(s)链接',
                'type': 'string',
            }
        },
        'required': ['url'],
    }

    def __init__(self, cfg: Optional[Dict] = None):
        super().__init__(cfg)
        self.data_root = self.cfg.get('path', os.path.join(DEFAULT_WORKSPACE, 'tools', self.name))
        self.extract_image = self.cfg.get('extract_image', False)
        self.structured_doc = self.cfg.get('structured_doc', False)

        self.db = Storage({'storage_root_path': self.data_root})

    def call(self, params: Union[str, dict], **kwargs) -> Union[str, list]:
        """Parse pdf by url, and return the formatted content.

        Returns:
            Extracted doc as plain text or the following list format:
              [
                {'page_num': 1,
                'content': [
                              {'text': 'This is one paragraph'},
                              {'table': 'This is one table'}
                           ],
                'title': 'If extracted, this is the title of the doc.'},
                {'page_num': 2,
                'content': [
                              {'text': 'This is one paragraph'},
                              {'table': 'This is one table'}
                           ]}
              ]
        """

        params = self._verify_json_format_args(params)
        path = params['url']
        cached_name_ori = f'{hash_sha256(path)}_ori'
        try:
            # Directly load the parsed doc
            parsed_file = self.db.get(cached_name_ori)
            parsed_file = json.loads(parsed_file)
            logger.info(f'Read parsed {path} from cache.')
        except KeyNotExistsError:
            logger.info(f'Start parsing {path}...')
            time1 = time.time()

            f_type = get_file_type(path)
            if f_type in PARSER_SUPPORTED_FILE_TYPES:
                if path.startswith('https://') or path.startswith('http://') or re.match(
                        r'^[A-Za-z]:\\', path) or re.match(r'^[A-Za-z]:/', path):
                    path = path
                else:
                    path = sanitize_chrome_file_path(path)

            os.makedirs(self.data_root, exist_ok=True)
            if is_http_url(path):
                # download online url
                tmp_file_root = os.path.join(self.data_root, hash_sha256(path))
                os.makedirs(tmp_file_root, exist_ok=True)
                path = save_url_to_local_work_dir(path, tmp_file_root)
            try:
                if f_type == 'pdf':
                    parsed_file = parse_pdf(path, self.extract_image)
                elif f_type == 'docx':
                    parsed_file = parse_word(path, self.extract_image)
                elif f_type == 'pptx':
                    parsed_file = parse_ppt(path, self.extract_image)
                elif f_type == 'txt':
                    parsed_file = parse_txt(path)
                elif f_type == 'html':
                    parsed_file = parse_html_bs(path, self.extract_image)
                elif f_type == 'csv':
                    parsed_file = parse_csv(path, self.extract_image)
                elif f_type == 'tsv':
                    parsed_file = parse_tsv(path, self.extract_image)
                elif f_type in ['xlsx', 'xls']:
                    parsed_file = parse_excel(path, self.extract_image)
                else:
                    _t = '/'.join(PARSER_SUPPORTED_FILE_TYPES)
                    raise ValueError(
                        f'Failed: The current parser does not support this file type! Supported types: {_t}')
            except Exception as ex:
                exception_type = type(ex).__name__
                exception_message = str(ex)
                raise DocParserError(code=exception_type, message=exception_message)

            for page in parsed_file:
                for para in page['content']:
                    # Todo: More attribute types
                    para['token'] = count_tokens(para.get('text', para.get('table')))
            time2 = time.time()
            logger.info(f'Finished parsing {path}. Time spent: {time2 - time1} seconds.')
            # Cache the parsing doc
            self.db.put(cached_name_ori, json.dumps(parsed_file, ensure_ascii=False, indent=2))

        if not self.structured_doc:
            return get_plain_doc(parsed_file)
        else:
            return parsed_file
