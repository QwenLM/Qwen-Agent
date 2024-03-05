import re

from qwen_agent.utils.tokenization_qwen import count_tokens

ONE_PAGE_TOKEN = 500


def rm_newlines(text):
    text = re.sub(r'(?<=[^\.。:：])\n', ' ', text)
    return text


def rm_cid(text):
    text = re.sub(r'\(cid:\d+\)', '', text)
    return text


def rm_hexadecimal(text):
    text = re.sub(r'[0-9A-Fa-f]{21,}', '', text)
    return text


def rm_continuous_placeholders(text):
    text = re.sub(r'(\.|-｜——｜。｜_|\*){7,}', '...', text)
    return text


def deal(text):
    text = rm_newlines(text)
    text = rm_cid(text)
    text = rm_hexadecimal(text)
    text = rm_continuous_placeholders(text)
    return text


def parse_doc(path):
    if '.pdf' in path.lower():
        from pdfminer.high_level import extract_text
        text = extract_text(path)
    elif '.docx' in path.lower():
        import docx2txt
        text = docx2txt.process(path)
    elif '.pptx' in path.lower():
        from pptx import Presentation
        ppt = Presentation(path)
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        text = '\n'.join(text)
    else:
        raise TypeError

    text = deal(text)
    return split_text_to_trunk(text, path)


def pre_process_html(s):
    # replace multiple newlines
    s = re.sub('\n+', '\n', s)
    # replace special string
    s = s.replace("Add to Qwen's Reading List", '')
    return s


def parse_html_bs(path):
    try:
        from bs4 import BeautifulSoup
    except Exception:
        raise ValueError('Please install bs4 by `pip install beautifulsoup4`')
    bs_kwargs = {'features': 'lxml'}
    with open(path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, **bs_kwargs)

    text = soup.get_text()

    if soup.title:
        title = str(soup.title.string)
    else:
        title = ''
    text = pre_process_html(text)
    return split_text_to_trunk(text, path, title)


def split_text_to_trunk(content: str, path: str, title: str = ''):
    all_tokens = count_tokens(content)
    all_pages = round(all_tokens / ONE_PAGE_TOKEN)
    if all_pages == 0:
        all_pages = 1
    len_content = len(content)
    len_one_page = int(len_content /
                       all_pages)  # Approximately equal to ONE_PAGE_TOKEN

    res = []
    for i in range(0, len_content, len_one_page):
        text = content[i:min(i + len_one_page, len_content)]
        res.append({
            'page_content': text,
            'metadata': {
                'source': path,
                'title': title,
                'page': (i % len_one_page)
            },
            'token': count_tokens(text)
        })
    return res
